"""Universal document extraction and text chunking."""

from __future__ import annotations

import csv
from pathlib import Path

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx", ".txt", ".csv"}


def extract_text(file_path: str) -> str:
    path = Path(file_path)
    suffix = path.suffix.lower()
    if suffix not in SUPPORTED_EXTENSIONS:
        display_suffix = suffix or "noma'lum"
        raise ValueError(f"Qo'llab-quvvatlanmaydigan format: {display_suffix}")
    if suffix == ".pdf":
        return _pdf(path)
    if suffix == ".docx":
        return _docx(path)
    if suffix == ".pptx":
        return _pptx(path)
    if suffix == ".xlsx":
        return _xlsx(path)
    if suffix == ".csv":
        return _csv(path)
    return path.read_text(encoding="utf-8-sig", errors="replace").strip()


def _pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    if reader.is_encrypted:
        raise ValueError("Bu PDF parol bilan himoyalangan.")
    pages = []
    for number, page in enumerate(reader.pages, 1):
        text = (page.extract_text() or "").strip()
        if text:
            pages.append(f"[SAHIFA {number}]\n{text}")
    result = "\n\n".join(pages).strip()
    if not result:
        raise ValueError("PDF ichidan matn topilmadi. Skanerlangan PDF uchun OCR kerak.")
    return result


def _docx(path: Path) -> str:
    doc = Document(str(path))
    parts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for index, table in enumerate(doc.tables, 1):
        rows = [" | ".join(cell.text.strip() for cell in row.cells) for row in table.rows]
        if rows:
            parts.append(f"[JADVAL {index}]\n" + "\n".join(rows))
    return "\n\n".join(parts).strip()


def _pptx(path: Path) -> str:
    presentation = Presentation(str(path))
    slides = []
    for number, slide in enumerate(presentation.slides, 1):
        texts = [shape.text.strip() for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        if texts:
            slides.append(f"[SLAYD {number}]\n" + "\n".join(texts))
    return "\n\n".join(slides).strip()


def _xlsx(path: Path) -> str:
    workbook = load_workbook(filename=str(path), read_only=True, data_only=True)
    sheets = []
    try:
        for sheet in workbook.worksheets:
            rows = []
            for row in sheet.iter_rows(values_only=True):
                values = [str(value).strip() for value in row if value is not None and str(value).strip()]
                if values:
                    rows.append(" | ".join(values))
            if rows:
                sheets.append(f"[VARAQ: {sheet.title}]\n" + "\n".join(rows))
    finally:
        workbook.close()
    return "\n\n".join(sheets).strip()


def _csv(path: Path) -> str:
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        rows = [" | ".join(cell.strip() for cell in row) for row in csv.reader(handle)]
    return "\n".join(row for row in rows if row).strip()


def split_into_chunks(text: str, max_words: int = 300, min_words: int = 50) -> list[str]:
    if not text or max_words <= 0 or min_words < 0:
        return []
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    chunks, current, count = [], [], 0
    for paragraph in paragraphs:
        words = paragraph.split()
        if len(words) > max_words:
            if current:
                chunks.append("\n\n".join(current))
                current, count = [], 0
            chunks.extend(" ".join(words[i:i + max_words]) for i in range(0, len(words), max_words))
            continue
        if current and count + len(words) > max_words and count >= min_words:
            chunks.append("\n\n".join(current))
            current, count = [], 0
        current.append(paragraph)
        count += len(words)
    if current:
        chunks.append("\n\n".join(current))
    return [chunk.strip() for chunk in chunks if chunk.strip()]
